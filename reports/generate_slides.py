from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path
import pandas as pd

# Paths
BASE_DIR = Path(__file__).resolve().parent.parent
PROCESSED_DIR = BASE_DIR / "data" / "processed"
REPORTS_DIR = BASE_DIR / "reports"

# Load data
metrics = pd.read_csv(PROCESSED_DIR / "performance_metrics.csv")
fund_master = pd.read_csv(PROCESSED_DIR / "01_fund_master.csv")
top5 = metrics.sort_values("cagr_pct", ascending=False).head(5)

# Create presentation
prs = Presentation()

# ── Slide 1 — Title ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Bluestock Mutual Fund Analytics"
slide.placeholders[1].text = "Capstone Project — Data Analyst Internship\nBluestock Fintech"

# ── Slide 2 — Project Overview ───────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Project Overview"
tf = slide.placeholders[1].text_frame
tf.text = "Objective"
p = tf.add_paragraph()
p.text = "• Analyze 40 mutual fund schemes across 10 fund houses"
p = tf.add_paragraph()
p.text = "• Build ETL pipeline, EDA, performance metrics and dashboard"
p = tf.add_paragraph()
p.text = "• Dataset: 46,000 NAV records, 32,778 investor transactions"
p = tf.add_paragraph()
p.text = "• Tools: Python, Pandas, SQLite, Streamlit"

# ── Slide 3 — Dataset Summary ────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Dataset Summary"
tf = slide.placeholders[1].text_frame
tf.text = "10 CSV files provided:"
datasets = [
    "• fund_master — 40 schemes, 15 columns",
    "• nav_history — 46,000 rows of daily NAV",
    "• investor_transactions — 32,778 records",
    "• benchmark_indices — NIFTY50, SENSEX data",
    "• portfolio_holdings, AUM, SIP inflows and more"
]
for d in datasets:
    p = tf.add_paragraph()
    p.text = d

# ── Slide 4 — Top 5 Funds ────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Top 5 Funds by CAGR"
tf = slide.placeholders[1].text_frame
tf.text = "Best performing funds:"
for _, row in top5.iterrows():
    name = str(row["scheme_name"])[:45] if pd.notna(row["scheme_name"]) else str(row["amfi_code"])
    p = tf.add_paragraph()
    p.text = f"• {name} — CAGR: {row['cagr_pct']}%"

# ── Slide 5 — Key Findings ───────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Findings"
tf = slide.placeholders[1].text_frame
tf.text = "Insights from analysis:"
findings = [
    "• Top fund achieved 31.48% CAGR over 4 years",
    "• Small cap funds outperformed large cap funds",
    "• Most funds have Sharpe Ratio above 1.0",
    "• SIP inflows grew consistently year on year",
    "• Mid cap funds offer best risk-adjusted returns",
]
for f in findings:
    p = tf.add_paragraph()
    p.text = f

# ── Slide 6 — Dashboard ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Interactive Dashboard"
tf = slide.placeholders[1].text_frame
tf.text = "Built with Streamlit:"
features = [
    "• Page 1 — Fund Overview with filters",
    "• Page 2 — NAV Growth over time",
    "• Page 3 — Performance Metrics table",
    "• Page 4 — Top 10 performers chart",
    "• Bonus B2 achieved — Streamlit as Power BI alternative",
]
for f in features:
    p = tf.add_paragraph()
    p.text = f

# ── Slide 7 — Conclusion ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
tf = slide.placeholders[1].text_frame
tf.text = "Summary:"
conclusions = [
    "• Successfully built end-to-end mutual fund analytics platform",
    "• ETL pipeline fetches and processes live NAV data",
    "• SQLite database stores all 9 tables",
    "• Fund recommender helps investors choose right funds",
    "• All deliverables D1-D7 completed",
]
for c in conclusions:
    p = tf.add_paragraph()
    p.text = c

# Save
prs.save(str(REPORTS_DIR / "Presentation.pptx"))
print("✅ Presentation.pptx created!")